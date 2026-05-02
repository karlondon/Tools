"""
GCP Engineering Team Journey — 2-slide Town Hall deck
Fresh presentation (no template), dark background, Google colours.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Fresh 16:9 presentation ───────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]  # truly blank

# ── Google Brand Palette ──────────────────────────────────────────────────────
GB   = RGBColor(0x42, 0x85, 0xF4)
GR   = RGBColor(0xEA, 0x43, 0x35)
GY   = RGBColor(0xFB, 0xBC, 0x04)
GG   = RGBColor(0x34, 0xA8, 0x53)
DARK = RGBColor(0x0D, 0x11, 0x17)
CARD = RGBColor(0x16, 0x1E, 0x2A)
C2   = RGBColor(0x1C, 0x26, 0x34)
WHT  = RGBColor(0xFF, 0xFF, 0xFF)
MUT  = RGBColor(0x8B, 0x96, 0xA8)
BDR  = RGBColor(0x28, 0x34, 0x44)
CYN  = RGBColor(0x06, 0xB6, 0xD4)
PUR  = RGBColor(0xA8, 0x55, 0xF7)
AMB  = RGBColor(0xF5, 0x9E, 0x0B)

# ── Primitives ────────────────────────────────────────────────────────────────
def bx(sl, l, t, w, h, fill=None, stroke=None, sw=0.75):
    s = sl.shapes.add_shape(1, l, t, w, h)
    if fill:   s.fill.solid(); s.fill.fore_color.rgb = fill
    else:      s.fill.background()
    if stroke: s.line.color.rgb = stroke; s.line.width = Pt(sw)
    else:      s.line.fill.background()
    return s

def tx(sl, text, l, t, w, h, sz=12, bold=False, color=WHT,
       align=PP_ALIGN.LEFT, font="Calibri"):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb

def ell(sl, cx, cy, r, fill):
    s = sl.shapes.add_shape(9, cx-r, cy-r, r*2, r*2)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def slide_bg(sl):
    bx(sl, 0, 0, SW, SH, fill=DARK)

def gbar(sl):
    for j, c in enumerate([GB, GR, GY, GG]):
        bx(sl, Inches(j*3.3325), 0, Inches(3.3325), Pt(6), fill=c)

def ftr(sl, n):
    bx(sl, 0, SH-Pt(2), SW, Pt(2), fill=BDR)
    tx(sl, "GCP Engineering  ·  Team Town Hall",
       Inches(0.4), SH-Inches(0.32), Inches(8), Inches(0.28), sz=8, color=MUT)
    tx(sl, f"0{n} / 02", SW-Inches(1.1), SH-Inches(0.32),
       Inches(0.9), Inches(0.28), sz=8, color=MUT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — HERO + VERTICAL TIMELINE
# Left 38%: big headline, narrative, 4 stat tiles
# Right 62%: 8-node timeline
# ════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
slide_bg(s1)
gbar(s1)

# ── LEFT ─────────────────────────────────────────────────────────────────────
LW = Inches(4.8)

tx(s1, "CMS  ·  GCP ENGINEERING PRACTICE",
   Inches(0.4), Inches(0.24), LW, Inches(0.3), sz=9, color=GB)

tx(s1, "From Zero\nto Full Orbit",
   Inches(0.4), Inches(0.6), LW, Inches(2.1),
   sz=52, bold=True, color=WHT, font="Calibri Light")

# Google 4-colour underline
for j, c in enumerate([GB, GR, GY, GG]):
    bx(s1, Inches(0.4 + j*0.46), Inches(2.76), Inches(0.42), Pt(5), fill=c)

tx(s1, "A team that inherited nothing, built everything,\n"
       "and is now a full GCP Engineering practice\n"
       "— in under a year.",
   Inches(0.4), Inches(2.94), Inches(4.4), Inches(0.88), sz=12, color=MUT)

# Date pill
bx(s1, Inches(0.4), Inches(3.96), Inches(3.6), Inches(0.42),
   fill=RGBColor(0x08,0x14,0x28), stroke=GB, sw=1.2)
tx(s1, "🚀  June 2025  →  Present  ·  Ongoing",
   Inches(0.4), Inches(3.98), Inches(3.6), Inches(0.38),
   sz=10, color=GB, align=PP_ALIGN.CENTER)

# 4 stat tiles
stats = [("11","MONTHS",GB),("2→3+","TEAM",GG),("4+","CLIENTS",GY),("1st","REVENUE",GR)]
TW = Inches(1.08); TH = Inches(1.18); TX0 = Inches(0.4); TY0 = Inches(4.56)
for k,(num,lbl,c) in enumerate(stats):
    x = TX0 + k*(TW+Inches(0.12))
    bx(s1, x, TY0, TW, TH, fill=CARD, stroke=c, sw=1.0)
    bx(s1, x, TY0, TW, Pt(5), fill=c)
    tx(s1, num, x, TY0+Inches(0.1), TW, Inches(0.6),
       sz=26, bold=True, color=c, align=PP_ALIGN.CENTER)
    tx(s1, lbl, x, TY0+Inches(0.72), TW, Inches(0.3),
       sz=8, color=MUT, align=PP_ALIGN.CENTER)

# ── RIGHT — Timeline panel ────────────────────────────────────────────────────
RX = Inches(5.18)
RW = SW - RX - Inches(0.3)

bx(s1, RX, 0, RW+Inches(0.3), SH, fill=CARD)
bx(s1, RX, 0, Pt(3), SH, fill=GB)

tx(s1, "THE JOURNEY", RX+Inches(0.3), Inches(0.22), RW, Inches(0.3),
   sz=9, bold=True, color=GB)
tx(s1, "8 milestones  ·  Jun 2025 → May 2026",
   RX+Inches(0.3), Inches(0.52), RW, Inches(0.3), sz=10, color=MUT)

nodes = [
    ("Jun 2025",     "🔄 Service Transition",
     "Handover — knowledge transfer & risk documentation",  GR),
    ("Jul 2025",     "🌱 Team Established",
     "2 people, 2 streams — Sr Manager + Consultant",       GB),
    ("Aug–Nov 2025", "🏗️  Vending Pipeline Delivered",
     "GCP Project Vending Pipeline scoped, built & live",   GG),
    ("Nov 2025",     "⭐ First Chargeable Win",
     "West Brom — Storage Transfer Service, AWS→GCP",       GY),
    ("Jan 2026",     "🧪 Live Validation",
     "GC Resell — pipeline proven in production",           GB),
    ("Feb 2026",     "👤 3rd Team Member Joins",
     "First dedicated hire as demand accelerates",          GG),
    ("Mar 2026",     "📣 Hiring for Scale",
     "3 open roles to meet incoming client demand",         PUR),
    ("Apr/May 2026", "🌐 YHCR Onboards",
     "New client extending existing GCP Org",               CYN),
]

SPX   = RX + Inches(0.52)
SY    = Inches(0.92)
NGAP  = Inches(0.8)
DR    = Inches(0.10)

bx(s1, SPX-Pt(1), SY, Pt(2), NGAP*(len(nodes)-1)+DR*2, fill=BDR)

for i,(date,title,body,c) in enumerate(nodes):
    ny = SY + i*NGAP
    ell(s1, SPX, ny+DR, DR+Inches(0.035), fill=RGBColor(c[0]//7,c[1]//7,c[2]//7))
    ell(s1, SPX, ny+DR, DR, fill=c)
    CX2 = SPX+Inches(0.26); CW2 = RW-Inches(0.46)
    bx(s1, CX2, ny-Inches(0.07), CW2, Inches(0.76), fill=C2, stroke=BDR, sw=0.5)
    bx(s1, CX2, ny-Inches(0.07), Inches(0.05), Inches(0.76), fill=c)
    tx(s1, date, CX2+Inches(0.1), ny-Inches(0.05),
       Inches(1.3), Inches(0.24), sz=8, bold=True, color=c)
    tx(s1, title, CX2+Inches(0.1), ny+Inches(0.18),
       CW2-Inches(0.18), Inches(0.24), sz=10, bold=True, color=WHT)
    tx(s1, body, CX2+Inches(0.1), ny+Inches(0.4),
       CW2-Inches(0.18), Inches(0.26), sz=8.5, color=MUT)

ftr(s1, 1)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — TEAM · CLIENTS · MILESTONES · FUTURE
# Three clean horizontal bands
# ════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
slide_bg(s2)
gbar(s2)

tx(s2, "GCP Engineering  ·  People  ·  Clients  ·  Milestones  ·  What's Next",
   Inches(0.4), Inches(0.18), Inches(12), Inches(0.28), sz=9, color=MUT)

# ── BAND A (top, y=0.55–2.72) ─────────────────────────────────────────────────
# Left half: Team Growth   Right half: Client Portfolio
BAND_A_Y = Inches(0.55)
DIV_X    = Inches(6.6)

# Section labels
tx(s2, "TEAM GROWTH ARC",
   Inches(0.4), BAND_A_Y, Inches(5.5), Inches(0.28), sz=9, bold=True, color=GG)
bx(s2, Inches(0.4), BAND_A_Y+Inches(0.28), Inches(5.5), Pt(2), fill=GG)

tx(s2, "CLIENT PORTFOLIO",
   DIV_X+Inches(0.2), BAND_A_Y, Inches(6.4), Inches(0.28), sz=9, bold=True, color=GB)
bx(s2, DIV_X+Inches(0.2), BAND_A_Y+Inches(0.28), Inches(6.4), Pt(2), fill=GB)

# Team phases — 4 tiles
phases = [
    ("2",   "people\nJul 2025",  GB),
    ("2",   "streams\nshared",   GG),
    ("3",   "members\nFeb 2026", GY),
    ("5+",  "target\n2026+",     GR),
]
PW2=Inches(1.18); PH2=Inches(1.74)
PX2=Inches(0.4); PY2=BAND_A_Y+Inches(0.38); PG2=Inches(0.14)
for j,(num,lbl,c) in enumerate(phases):
    x = PX2+j*(PW2+PG2)
    bx(s2, x, PY2, PW2, PH2, fill=CARD, stroke=c, sw=1.0)
    bx(s2, x, PY2, PW2, Pt(5), fill=c)
    tx(s2, num, x, PY2+Inches(0.12), PW2, Inches(0.7),
       sz=32, bold=True, color=c, align=PP_ALIGN.CENTER)
    tx(s2, lbl, x, PY2+Inches(0.86), PW2, Inches(0.62),
       sz=9, color=MUT, align=PP_ALIGN.CENTER)
    if j < 3:
        tx(s2, "›", x+PW2, PY2+Inches(0.6), PG2+Inches(0.02),
           Inches(0.32), sz=16, bold=True, color=GG, align=PP_ALIGN.CENTER)

# Client cards — 2×2 grid
clients = [
    ("West Brom",           "Nov 2025",    "● Active",
     "Storage Transfer Service — AWS→GCP backup",  GG),
    ("GC Resell",           "Jan 2026",    "● Live",
     "Vending pipeline live in production",         GB),
    ("Alpha Commerce\n+PairD","Aug 2025",  "● Internal",
     "Billing exports & ops support",              GY),
    ("YHCR",                "Apr/May 2026","◐ Incoming",
     "Extending existing GCP Org",                  CYN),
]
CCW=Inches(3.0); CCH=Inches(0.82)
CCX0=DIV_X+Inches(0.2); CCY0=BAND_A_Y+Inches(0.38)
CCGX=Inches(0.22); CCGY=Inches(0.12)
for idx,(name,since,status,desc,c) in enumerate(clients):
    ci=idx%2; ri=idx//2
    x=CCX0+ci*(CCW+CCGX); y=CCY0+ri*(CCH+CCGY)
    bx(s2, x, y, CCW, CCH, fill=CARD, stroke=c, sw=0.8)
    bx(s2, x, y, Inches(0.055), CCH, fill=c)
    tx(s2, name, x+Inches(0.12), y+Inches(0.05),
       Inches(1.8), Inches(0.38), sz=10, bold=True, color=WHT)
    tx(s2, since, x+Inches(0.12), y+Inches(0.44),
       Inches(1.1), Inches(0.22), sz=8, color=MUT)
    tx(s2, status, x+CCW-Inches(1.1), y+Inches(0.07),
       Inches(1.0), Inches(0.22), sz=8, bold=True, color=c, align=PP_ALIGN.RIGHT)
    tx(s2, desc, x+Inches(0.12), y+Inches(0.44),
       CCW-Inches(0.18), Inches(0.34), sz=8.5, color=MUT)

# Vertical divider
bx(s2, DIV_X, BAND_A_Y, Pt(1), Inches(2.22), fill=BDR)

# ── BAND B — Milestone ribbon (y=2.9–3.9) ────────────────────────────────────
BAND_B_Y = Inches(2.9)
bx(s2, Inches(0.4), BAND_B_Y-Pt(1), SW-Inches(0.8), Pt(1), fill=BDR)
tx(s2, "KEY MILESTONES — Jun 2025 to May 2026",
   Inches(0.4), BAND_B_Y+Inches(0.06), Inches(6), Inches(0.26),
   sz=9, bold=True, color=MUT)

milestones = [
    ("🔄","Transition",  GR),
    ("🏗️","Pipeline",    GG),
    ("⭐", "West Brom",  GY),
    ("🧪","Validation",  GB),
    ("👤","3rd Hire",    GG),
    ("📣","3 Roles",     PUR),
    ("🌐","YHCR",        CYN),
]
MW2=Inches(1.64); MH2=Inches(0.96)
MX2=Inches(0.4); MY2=BAND_B_Y+Inches(0.36); MG2=Inches(0.15)
for j,(icon,label,c) in enumerate(milestones):
    x=MX2+j*(MW2+MG2)
    bx(s2, x, MY2, MW2, MH2, fill=c)
    tx(s2, icon, x, MY2+Inches(0.06), MW2, Inches(0.42),
       sz=22, align=PP_ALIGN.CENTER, color=WHT)
    tx(s2, label, x, MY2+Inches(0.54), MW2, Inches(0.34),
       sz=9, bold=True, color=WHT, align=PP_ALIGN.CENTER)
    if j<6:
        bx(s2, x+MW2, MY2+MH2/2-Pt(1), MG2, Pt(2),
           fill=RGBColor(0x3A,0x46,0x58))

# ── BAND C — What's Next (y=4.08 — fits within 6.9 safe area) ────────────────
BAND_C_Y = Inches(4.06)
bx(s2, Inches(0.4), BAND_C_Y-Pt(1), SW-Inches(0.8), Pt(1), fill=BDR)
tx(s2, "WHAT'S NEXT",
   Inches(0.4), BAND_C_Y+Inches(0.04), Inches(4), Inches(0.24),
   sz=9, bold=True, color=MUT)

# Statement block — height capped so bottom = 6.6"
SB_Y  = BAND_C_Y + Inches(0.34)
SB_H  = Inches(1.36)   # ends at 4.06+0.34+1.36 = 5.76"
bx(s2, Inches(0.4), SB_Y, Inches(4.2), SB_H, fill=C2, stroke=BDR, sw=0.75)
sh2 = SB_H / 4
for j,c in enumerate([GB,GR,GY,GG]):
    bx(s2, Inches(0.4), SB_Y+j*sh2, Inches(0.07), sh2, fill=c)
tx(s2, "The best chapters\nare still being\nwritten.",
   Inches(0.58), SB_Y+Inches(0.14), Inches(3.9), Inches(1.16),
   sz=22, bold=True, color=WHT, font="Calibri Light")

# 5 future pillars — same height as statement block
pillars = [
    ("👥","Team\nExpansion",  GB),
    ("🌐","New\nClients",     GG),
    ("⚙️", "Automation",      GY),
    ("📊","Revenue\nGrowth",  GR),
    ("☁️", "Cloud\nPractice", PUR),
]
PLW=Inches(1.56); PLH=SB_H
PLX=Inches(5.0); PLY=SB_Y; PLG=Inches(0.2)
for j,(icon,label,c) in enumerate(pillars):
    x=PLX+j*(PLW+PLG)
    bx(s2, x, PLY, PLW, PLH, fill=CARD, stroke=c, sw=1.0)
    bx(s2, x, PLY, PLW, Pt(5), fill=c)
    tx(s2, icon, x, PLY+Inches(0.1), PLW, Inches(0.48),
       sz=26, align=PP_ALIGN.CENTER, color=WHT)
    tx(s2, label, x, PLY+Inches(0.62), PLW, Inches(0.62),
       sz=9, bold=True, color=c, align=PP_ALIGN.CENTER)

ftr(s2, 2)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "GCP-Team-Journey/GCP_Engineering_Team_Journey.pptx"
prs.save(out)
print(f"Saved: {out}")